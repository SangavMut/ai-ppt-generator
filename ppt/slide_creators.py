import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class SlideCreators:
    def __init__(self, presentation: Presentation, image_downloader):
        self.presentation = presentation
        self.image_downloader = image_downloader  # ImageDownloader instance

    def create_title_slide(self, title, subtitle=""):
        """Create a title slide"""
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        return slide

    def create_content_slide(self, title, content, include_image=False):
        """Create a content slide with bullet points"""
        slide_layout = self.presentation.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True

        if isinstance(content, list):
            content = "\n".join(content)

        content_shape = slide.placeholders[1]
        content_shape.text = content

        text_frame = content_shape.text_frame
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(51, 51, 51)

        if include_image:
            image_path = self.image_downloader.download_image(content)  # Download image based on content
            slide.shapes.add_picture(image_path, Inches(6), Inches(2), height=Inches(4))

        return slide

    def create_image_slide(self, title, content, image_query):
        """Create a slide with image and text"""
        slide_layout = self.presentation.slide_layouts[8]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Add content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(5))
        content_frame = content_box.text_frame

        if isinstance(content, list):
            content = "\n".join(content)
        content_frame.text = content

        # Format content
        for paragraph in content_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(51, 51, 51)

        # Add image
        try:
            image_path = self.image_downloader.download_image(image_query)  # Get image based on query
            if image_path and os.path.exists(image_path):
                slide.shapes.add_picture(image_path, Inches(5), Inches(2), height=Inches(5))
                os.remove(image_path)  # Clean up temporary image file
        except Exception as e:
            print(f"Error adding image: {e}")

        return slide
